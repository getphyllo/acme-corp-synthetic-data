"""PPTX deck toolkit for the Acme report pack (decks 41-60).

Same contract as lib.Doc: a small set of fixed 16:9 layouts, one call per slide,
every number supplied by the builder from acme.duckdb / seeds. Charts are the
same PNGs that lib.chart_* produces for the PDFs, so a deck and a PDF on the
same subject look like they came out of the same template — because they did.

    k = Deck("41-....pptx", kicker="QUARTERLY BUSINESS REVIEW",
             title="Acme Corp — Q2 FY2026 Company Business Review",
             subtitle="Enterprise performance, all brands and channels",
             byline="Diane Halverson, VP Sales NA  ·  July 2026",
             short="Q2 2026 Company BR")
    k.agenda([...])
    k.exec_summary("lede ...", tiles=[("−5.3%", "Q1 revenue vs plan"), ...],
                   bullets=[...])
    k.chart("THE NUMBERS", "Revenue against plan", png, note="...")
    k.chart_table("SHARE", "Where share moved", png, headers, rows, note="...")
    k.reco("THE ASK", "Recommendations", [(action, owner, when), ...])
    k.close("One number to remember", ["..."])
    k.build()
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from lib import palette, OUT_DIR

# ------------------------------------------------------------------ colors ---
def _rgb(name):
    return RGBColor.from_string(palette[name].lstrip("#"))

NAVY, INK, TEAL, GOLD, RUST = (_rgb(k) for k in ("navy", "ink", "teal", "gold", "rust"))
SKY, SLATE, MIST, MIST2, LINE = (_rgb(k) for k in ("sky", "slate", "mist", "mist2", "line"))
WHITE, GREEN, AMBER = _rgb("white"), _rgb("green"), _rgb("amber")
FONT = "Calibri"

# 16:9 geometry (13.333 x 7.5 in)
SW, SH = Inches(13.333), Inches(7.5)
M = Inches(0.7)                      # left/right margin
CW = Inches(13.333 - 1.4)            # content width
TOP = Inches(1.55)                   # top of the content band
BOT = Inches(6.62)                   # bottom of the content band
FOOTER = "Internal · Confidential — Acme Corp synthetic demo data. Fictional; not for real decision-making."


def _lines(text, width_in, size):
    """Rough wrapped-line count. Calibrated: 7.1in at 14.5pt wraps at ~75 chars."""
    cpl = max(8, int(width_in * 153 / size))
    n = 0
    for ln in (text if isinstance(text, (list, tuple)) else [text]):
        n += max(1, -(-len(str(ln).replace("**", "")) // cpl))
    return n


def _est_h(text, width_in, size, leading=1.2, space_after=0.0):
    """Estimated rendered height, in EMU, of a text block."""
    n = len(text) if isinstance(text, (list, tuple)) else 1
    h = _lines(text, width_in, size) * size * leading / 72.0
    return Inches(h + space_after * n)


class Deck:
    def __init__(self, filename, title, kicker, subtitle, byline, short,
                 doc_type="Internal management presentation"):
        self.filename = filename
        self.path = os.path.join(OUT_DIR, filename)
        self.title, self.kicker, self.subtitle = title, kicker, subtitle
        self.byline, self.short, self.doc_type = byline, short, doc_type
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SW, SH
        self.prs.core_properties.title = title
        self.prs.core_properties.author = "Acme Corp (synthetic)"
        self.prs.core_properties.comments = doc_type
        self.n = 0
        self._title_slide()

    # ------------------------------------------------------------ primitives --
    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, s, l, t, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None):
        sp = s.shapes.add_shape(shape, l, t, w, h)
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(0.75)
        sp.text_frame.word_wrap = True
        return sp

    def _tb(self, s, l, t, w, h, text, size=12, color=INK, bold=False,
            align=PP_ALIGN.LEFT, space=4, caps=False, italic=False, leading=1.18):
        box = s.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        lines = text if isinstance(text, (list, tuple)) else [text]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space)
            p.line_spacing = leading
            self._runs(p, ln.upper() if caps else ln, size, color, bold, italic)
        return box

    def _runs(self, p, text, size, color, bold, italic=False):
        """Inline **bold** markup -> runs."""
        for i, chunk in enumerate(str(text).split("**")):
            if not chunk:
                continue
            r = p.add_run()
            r.text = chunk
            f = r.font
            f.name, f.size, f.color.rgb = FONT, Pt(size), color
            f.bold = bold or (i % 2 == 1)
            f.italic = italic
        if not p.runs:                                     # empty line keeps spacing
            r = p.add_run(); r.text = " "
            r.font.name, r.font.size = FONT, Pt(size)

    def _head(self, s, kicker, headline):
        self._tb(s, M, Inches(0.40), CW, Inches(0.28), kicker, size=10, color=TEAL,
                 bold=True, caps=True, space=0)
        self._tb(s, M, Inches(0.74), CW, Inches(0.62), headline, size=25, color=NAVY,
                 bold=True, space=0)
        ln = self._rect(s, M, Inches(1.36), Inches(1.1), Emu(22860), TEAL)
        ln.line.fill.background()

    def _foot(self, s):
        self.n += 1
        self._tb(s, M, Inches(7.02), Inches(9.2), Inches(0.24), FOOTER, size=7.5,
                 color=SLATE, space=0)
        self._tb(s, Inches(11.4), Inches(7.02), Inches(1.23), Inches(0.24),
                 f"{self.short}  ·  {self.n}", size=7.5, color=SLATE,
                 align=PP_ALIGN.RIGHT, space=0)

    def _slide(self, kicker, headline):
        s = self._blank()
        self._head(s, kicker, headline)
        self._foot(s)
        return s

    def _note(self, s, text, top=Inches(6.16), left=M, width=CW, color=SLATE, size=10.5):
        if text:
            self._tb(s, left, top, width, Inches(0.42), text, size=size, color=color,
                     space=0, leading=1.16)

    # ------------------------------------------------------------- title page --
    def _title_slide(self):
        s = self._blank()
        self._rect(s, 0, 0, SW, SH, NAVY)
        for i in range(3):
            d = self._rect(s, Inches(0.92 + i * 0.34), Inches(0.92), Inches(0.2),
                           Inches(0.2), TEAL if i == 0 else _rgb("sky") if i == 1 else GOLD,
                           shape=MSO_SHAPE.OVAL)
            d.line.fill.background()
        self._tb(s, Inches(0.9), Inches(2.0), Inches(11.4), Inches(0.4), self.kicker,
                 size=12.5, color=_rgb("teal"), bold=True, caps=True, space=0)
        th = _est_h(self.title, 11.5, 41, 1.08)
        self._tb(s, Inches(0.88), Inches(2.55), Inches(11.5), th, self.title,
                 size=41, color=WHITE, bold=True, space=0, leading=1.08)
        self._tb(s, Inches(0.9), Inches(2.55) + th + Inches(0.42), Inches(11.3),
                 Inches(0.9), self.subtitle, size=16, color=_rgb("line"),
                 space=0, leading=1.2)
        self._tb(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.8),
                 [self.byline, self.doc_type + " · Internal · Confidential — synthetic demo data"],
                 size=11, color=_rgb("slate"), space=3)

    # ----------------------------------------------------------------- slides --
    def agenda(self, items, kicker="AGENDA", headline="What this deck covers"):
        s = self._slide(kicker, headline)
        n = len(items)
        rows = (n + 1) // 2
        colw = Inches(5.85)
        for i, it in enumerate(items):
            col, row = i // rows, i % rows
            l = M + (colw + Inches(0.24)) * col
            t = TOP + Inches(0.06) + Inches(0.72) * row
            bx = self._rect(s, l, t, colw, Inches(0.6), MIST2)
            bx.line.color.rgb = LINE
            bx.line.width = Pt(0.5)
            self._tb(s, l + Inches(0.14), t + Inches(0.14), Inches(0.4), Inches(0.3),
                     f"{i+1:02d}", size=13, color=TEAL, bold=True, space=0)
            self._tb(s, l + Inches(0.62), t + Inches(0.15), colw - Inches(0.8),
                     Inches(0.32), it, size=12.5, color=NAVY, bold=True, space=0)
        return s

    def tiles(self, kicker, headline, tiles, body=None, note=None):
        """tiles: list of (value, label) or (value, label, sub). 3-6 cards."""
        s = self._slide(kicker, headline)
        n = len(tiles)
        gap = Inches(0.22)
        w = Emu(int((CW - gap * (n - 1)) / n))
        for i, tl in enumerate(tiles):
            val, lab = tl[0], tl[1]
            sub = tl[2] if len(tl) > 2 else None
            l = M + Emu(int(w + gap)) * i
            card = self._rect(s, l, TOP + Inches(0.1), w, Inches(1.72), MIST2)
            card.line.color.rgb = LINE
            card.line.width = Pt(0.5)
            bar = self._rect(s, l, TOP + Inches(0.1), Emu(45720), Inches(1.72), TEAL)
            bar.line.fill.background()
            self._tb(s, l + Inches(0.2), TOP + Inches(0.3), w - Inches(0.34),
                     Inches(0.6), str(val), size=27, color=NAVY, bold=True, space=0)
            self._tb(s, l + Inches(0.2), TOP + Inches(0.96), w - Inches(0.34),
                     Inches(0.3), lab, size=10, color=SLATE, bold=True, caps=True, space=0)
            if sub:
                self._tb(s, l + Inches(0.2), TOP + Inches(1.26), w - Inches(0.34),
                         Inches(0.42), sub, size=9.5, color=SLATE, space=0, leading=1.1)
        if body:
            self._tb(s, M, TOP + Inches(2.1), CW, Inches(1.9),
                     body if isinstance(body, (list, tuple)) else [body],
                     size=13, color=INK, space=8, leading=1.24)
        self._note(s, note)
        return s

    def exec_summary(self, lede, tiles=None, bullets=None, kicker="THE HEADLINE",
                     headline="Executive summary"):
        s = self._slide(kicker, headline)
        left_w = Inches(7.1) if tiles else CW
        lede_h = _est_h(lede, left_w / 914400, 14.5, 1.26)
        self._tb(s, M, TOP, left_w, lede_h, lede, size=14.5, color=INK,
                 space=10, leading=1.26)
        if bullets:
            self._bullet_block(s, M, TOP + lede_h + Inches(0.24), left_w, bullets, size=12)
        if tiles:
            l0 = M + Inches(7.55)
            w = Inches(2.35)
            for i, tl in enumerate(tiles[:4]):
                col, row = i % 2, i // 2
                l = l0 + (w + Inches(0.24)) * col
                t = TOP + Inches(0.02) + Inches(1.42) * row
                card = self._rect(s, l, t, w, Inches(1.24), MIST2)
                card.line.color.rgb = LINE
                card.line.width = Pt(0.5)
                self._tb(s, l + Inches(0.16), t + Inches(0.16), w - Inches(0.3),
                         Inches(0.5), str(tl[0]), size=22, color=NAVY, bold=True, space=0)
                self._tb(s, l + Inches(0.16), t + Inches(0.72), w - Inches(0.3),
                         Inches(0.44), tl[1], size=9.5, color=SLATE, bold=True,
                         caps=True, space=0, leading=1.12)
            if len(tiles) > 4:
                extra = " · ".join(f"{t[0]} {t[1]}" for t in tiles[4:])
                self._tb(s, l0, TOP + Inches(2.92), Inches(4.94), Inches(0.6),
                         extra, size=10, color=SLATE, space=0, leading=1.15)
        return s

    def _bullet_block(self, s, l, t, w, bullets, size=12.5, gap=0.06):
        box = s.shapes.add_textbox(l, t, w, BOT - t)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(size * 0.62)
            p.line_spacing = 1.2
            r = p.add_run(); r.text = "—  "
            r.font.name, r.font.size, r.font.color.rgb, r.font.bold = FONT, Pt(size), TEAL, True
            self._runs(p, b, size, INK, False)
        return box

    def bullets(self, kicker, headline, bullets, lede=None, note=None, tiles=None):
        s = self._slide(kicker, headline)
        t = TOP
        w = CW if not tiles else Inches(7.9)
        if lede:
            lh = _est_h(lede, w / 914400, 14, 1.24)
            self._tb(s, M, t, w, lh, lede, size=14, color=SLATE,
                     italic=True, space=0, leading=1.24)
            t = t + lh + Inches(0.22)
        self._bullet_block(s, M, t, w, bullets, size=13)
        if tiles:
            l = M + Inches(8.3)
            for i, tl in enumerate(tiles[:4]):
                tt = TOP + Inches(1.28) * i
                card = self._rect(s, l, tt, Inches(3.53), Inches(1.1), MIST2)
                card.line.color.rgb = LINE
                card.line.width = Pt(0.5)
                self._tb(s, l + Inches(0.16), tt + Inches(0.13), Inches(3.2),
                         Inches(0.44), str(tl[0]), size=19, color=NAVY, bold=True, space=0)
                self._tb(s, l + Inches(0.16), tt + Inches(0.62), Inches(3.2),
                         Inches(0.4), tl[1], size=9.5, color=SLATE, bold=True,
                         caps=True, space=0, leading=1.12)
        self._note(s, note)
        return s

    def two_col(self, kicker, headline, left_title, left_items, right_title,
                right_items, note=None):
        s = self._slide(kicker, headline)
        w = Inches(5.85)
        for i, (ti, items) in enumerate(((left_title, left_items), (right_title, right_items))):
            l = M + (w + Inches(0.24)) * i
            hdr = self._rect(s, l, TOP, w, Inches(0.44), NAVY if i == 0 else _rgb("slate"))
            hdr.line.fill.background()
            self._tb(s, l + Inches(0.16), TOP + Inches(0.11), w - Inches(0.3),
                     Inches(0.3), ti, size=11.5, color=WHITE, bold=True, caps=True, space=0)
            self._bullet_block(s, l + Inches(0.04), TOP + Inches(0.62), w - Inches(0.1),
                               items, size=12)
        self._note(s, note)
        return s

    # ------------------------------------------------------------------ charts --
    def _fit(self, s, png, l, t, w, h):
        """Scale-to-fit a PNG inside the box (l,t,w,h) and centre it there."""
        pic = s.shapes.add_picture(png, l, t, width=w)
        scale = min(1.0, h / pic.height)
        if scale < 1.0:
            pic.width, pic.height = Emu(int(pic.width * scale)), Emu(int(pic.height * scale))
        pic.left = Emu(int(l + (w - pic.width) / 2))
        pic.top = Emu(int(t + (h - pic.height) / 2))
        return pic

    def chart(self, kicker, headline, png, note=None, lede=None, width=11.0):
        s = self._slide(kicker, headline)
        t = TOP
        if lede:
            lh = _est_h(lede, 11.9, 13, 1.2)
            self._tb(s, M, t, CW, lh, lede, size=13, color=SLATE,
                     italic=True, space=0, leading=1.2)
            t = t + lh + Inches(0.14)
        self._fit(s, png, Emu(int((SW - Inches(width)) / 2)), t, Inches(width),
                  Inches(6.05) - t)
        self._note(s, note, top=Inches(6.2))
        return s

    def charts2(self, kicker, headline, png1, png2, note=None, captions=None):
        s = self._slide(kicker, headline)
        w, h = Inches(5.9), Inches(3.85 if captions else 4.3)
        for i, png in enumerate((png1, png2)):
            l = M + (w + Inches(0.24)) * i
            pic = self._fit(s, png, l, TOP + Inches(0.06), w, h)
            if captions:
                self._tb(s, l, Emu(int(pic.top + pic.height + Inches(0.16))), w,
                         Inches(0.5), captions[i], size=10, color=SLATE,
                         space=0, leading=1.15)
        self._note(s, note, top=Inches(6.2))
        return s

    def chart_table(self, kicker, headline, png, headers, rows, note=None,
                    widths=None, chart_left=True, total_row=False, status_col=None,
                    align_right_from=1, callout=None, size=9.5):
        s = self._slide(kicker, headline)
        cw_chart, cw_tbl = Inches(6.9), Inches(4.83)
        lc = M if chart_left else M + cw_tbl + Inches(0.2)
        lt = M + cw_chart + Inches(0.2) if chart_left else M
        self._fit(s, png, lc, TOP + Inches(0.08), cw_chart, Inches(4.1))
        shp = self._table(s, lt, TOP + Inches(0.08), cw_tbl, headers, rows, widths,
                          total_row=total_row, size=size, status_col=status_col,
                          align_right_from=align_right_from, fill_h=Inches(4.1))
        note_top = Inches(6.2)
        if callout:
            box = self.callout(s, callout[0], callout[1],
                               callout[2] if len(callout) > 2 else "info",
                               top=Inches(5.4), left=lt, width=cw_tbl)
            note_top = self._after(box)
        self._note(s, note, top=note_top)
        return s

    def chart_bullets(self, kicker, headline, png, bullets, note=None, chart_left=True,
                      callout=None):
        s = self._slide(kicker, headline)
        cw_chart, cw_txt = Inches(7.1), Inches(4.6)
        lc = M if chart_left else M + cw_txt + Inches(0.25)
        lt = M + cw_chart + Inches(0.25) if chart_left else M
        self._fit(s, png, lc, TOP + Inches(0.06), cw_chart, Inches(4.2))
        self._bullet_block(s, lt, TOP + Inches(0.06), cw_txt, bullets, size=11.5)
        note_top = Inches(6.2)
        if callout:
            box = self.callout(s, callout[0], callout[1],
                               callout[2] if len(callout) > 2 else "info",
                               top=Inches(5.4), left=lt, width=cw_txt)
            note_top = self._after(box)
        self._note(s, note, top=note_top)
        return s

    # ------------------------------------------------------------------ tables --
    def _table(self, s, l, t, w, headers, rows, widths=None, total_row=False,
               size=10, highlight=None, status_col=None, align_right_from=1,
               fill_h=None):
        nr, nc = len(rows) + 1, len(headers)
        row_h = Inches(0.29 if nr > 9 else 0.33)
        if fill_h:
            # spread short tables over the band; shrink (to a floor) rather than overflow
            row_h = Emu(max(int(Inches(0.235)), min(int(Inches(0.5)), int(fill_h / nr))))
            if row_h * nr > fill_h:
                raise ValueError(
                    f"table of {nr} rows cannot fit {fill_h / 914400:.2f}in on one slide "
                    f"(needs {(row_h * nr) / 914400:.2f}in) — split it or cut rows")
        shp = s.shapes.add_table(nr, nc, l, t, w, row_h * nr)
        tbl = shp.table
        tbl.first_row = True
        if widths:
            tot = sum(widths)
            for i, fr in enumerate(widths):
                tbl.columns[i].width = Emu(int(w * fr / tot))
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = str(h)
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
            c.margin_left = c.margin_right = Inches(0.07)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j >= align_right_from else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(size - 0.5), True, WHITE
        for i, row in enumerate(rows, start=1):
            last = total_row and i == len(rows)
            for j, v in enumerate(row):
                c = tbl.cell(i, j)
                c.text = "" if v is None else str(v)
                c.fill.solid()
                c.fill.fore_color.rgb = MIST if last else (WHITE if i % 2 else MIST2)
                c.margin_left = c.margin_right = Inches(0.07)
                c.margin_top = c.margin_bottom = Inches(0.02)
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = c.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT if j >= align_right_from else PP_ALIGN.LEFT
                col = INK
                if status_col is not None and j == status_col:
                    col = {"On track": GREEN, "Watch": AMBER, "Action": RUST,
                           "Reinvest": GREEN, "Reallocate": RUST, "Hold": SLATE,
                           "Go": GREEN, "No-go": RUST, "Pass": GREEN, "Fail": RUST}.get(str(v), INK)
                for r in p.runs:
                    r.font.name, r.font.size, r.font.color.rgb = FONT, Pt(size), col
                    r.font.bold = bool(last or j == highlight or
                                       (status_col is not None and j == status_col))
        return shp

    def table(self, kicker, headline, headers, rows, note=None, widths=None,
              total_row=False, lede=None, status_col=None, highlight=None, size=10.5,
              align_right_from=1, callout=None, fill=True):
        s = self._slide(kicker, headline)
        t = TOP
        if lede:
            lh = _est_h(lede, 11.9, 13, 1.2)
            self._tb(s, M, t, CW, lh, lede, size=13, color=SLATE,
                     italic=True, space=0, leading=1.2)
            t = t + lh + Inches(0.16)
        band = Inches(5.3) - t if callout else Inches(6.05) - t
        shp = self._table(s, M, t, CW, headers, rows, widths, total_row=total_row,
                          size=size, status_col=status_col, highlight=highlight,
                          align_right_from=align_right_from,
                          fill_h=band if fill else None)
        note_top = Inches(6.2)
        if callout:
            box = self.callout(s, callout[0], callout[1],
                               callout[2] if len(callout) > 2 else "info",
                               top=Emu(int(shp.top + shp.height + Inches(0.28))))
            note_top = self._after(box)
        self._note(s, note, top=note_top)
        return s

    def callout(self, s, title, text, kind="info", top=Inches(5.5), left=M, width=None):
        """Colored-bar callout, placed on an existing slide. Height follows the text."""
        col = {"info": TEAL, "risk": RUST, "action": GOLD, "win": GREEN}.get(kind, TEAL)
        w = width or CW
        body_h = _est_h(text, w / 914400 - 0.4, 10.5, 1.16)
        h = Inches(0.5) + body_h
        box = self._rect(s, left, top, w, h, MIST2)
        box.line.fill.background()
        bar = self._rect(s, left, top, Emu(45720), h, col)
        bar.line.fill.background()
        self._tb(s, left + Inches(0.18), top + Inches(0.1), w - Inches(0.34),
                 Inches(0.26), title, size=11.5, color=NAVY, bold=True, space=0)
        self._tb(s, left + Inches(0.18), top + Inches(0.38), w - Inches(0.34),
                 body_h, text, size=10.5, color=INK, space=0, leading=1.16)
        return box

    def _after(self, box, fallback=Inches(6.2)):
        """Note position that clears a callout instead of colliding with it."""
        t = Emu(int(box.top + box.height + Inches(0.12)))
        return max(fallback, min(t, Inches(6.62)))

    def risk(self, kicker, headline, rows, note=None, lede=None, callout=None):
        """rows: (risk, impact, mitigation, owner)"""
        return self.table(kicker, headline, ["Risk", "Impact if it lands", "Mitigation", "Owner"],
                          rows, widths=[0.24, 0.24, 0.36, 0.16], note=note, lede=lede,
                          size=10, align_right_from=9, callout=callout)

    def reco(self, kicker, headline, items, note=None, lede=None, callout=None):
        """items: (action, owner, when)"""
        rows = [[str(i + 1), a, o, w] for i, (a, o, w) in enumerate(items)]
        return self.table(kicker, headline, ["#", "Action", "Owner", "When"], rows,
                          widths=[0.05, 0.55, 0.23, 0.17], note=note, lede=lede,
                          size=10.5, align_right_from=3, callout=callout)

    def close(self, headline, lines, kicker="THE ONE THING"):
        s = self._blank()
        self._rect(s, 0, 0, SW, SH, NAVY)
        bar = self._rect(s, Inches(0.9), Inches(1.5), Inches(1.1), Emu(27432), TEAL)
        bar.line.fill.background()
        self._tb(s, Inches(0.9), Inches(1.0), Inches(11.4), Inches(0.32), kicker,
                 size=12, color=TEAL, bold=True, caps=True, space=0)
        hh = _est_h(headline, 11.2, 32, 1.12)
        self._tb(s, Inches(0.88), Inches(1.95), Inches(11.2), hh, headline,
                 size=32, color=WHITE, bold=True, space=0, leading=1.12)
        self._tb(s, Inches(0.9), Inches(1.95) + hh + Inches(0.5), Inches(10.8), Inches(2.2),
                 lines if isinstance(lines, (list, tuple)) else [lines],
                 size=14, color=_rgb("line"), space=9, leading=1.24)
        self._tb(s, Inches(0.9), Inches(7.02), Inches(11.4), Inches(0.24),
                 FOOTER, size=7.5, color=SLATE, space=0)
        return s

    def build(self):
        self.prs.save(self.path)
        return self.path
